# core/generator.py

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import json # Import json for pretty printing

def create_ppt_from_template(slide_data, output_path, template_path=None):
    """
    Creates a NEW, separate PPT file from slide_data, applying styles from a template.
    This version includes debugging print statements to diagnose issues.
    """
    print("\n--- Starting PPT Generation ---")
    
    if template_path:
        print(f"Loading template from: {template_path}")
        prs = Presentation(template_path)
        
        # --- NEW, MORE RELIABLE SLIDE DELETION LOGIC ---
        print(f"Template has {len(prs.slides)} slides. Deleting them now.")
        
        # This is a safe way to remove all slides
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
        
        print("All original slides deleted successfully.")
        # --- END OF NEW LOGIC ---

    else:
        print("No template provided. Creating a blank presentation.")
        prs = Presentation()

    print("\n--- DEBUG: Incoming Slide Data ---")
    print(json.dumps(slide_data, indent=2))
    print("----------------------------------\n")

    # Find a suitable "Title and Content" layout
    title_and_content_layout = None
    try:
        title_and_content_layout = prs.slide_layouts[1]
        print("Found 'Title and Content' layout at index 1.")
    except IndexError:
        print("Warning: Layout at index 1 not found. Searching for a suitable layout...")
        for i, layout in enumerate(prs.slide_layouts):
            has_title = any(ph.placeholder_format.type == PP_PLACEHOLDER.TITLE for ph in layout.placeholders)
            has_body = any(ph.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) for ph in layout.placeholders)
            if has_title and has_body:
                title_and_content_layout = layout
                print(f"Found suitable layout '{layout.name}' at index {i}.")
                break
    
    if not title_and_content_layout:
        print("Error: No suitable 'Title and Content' layout found. Using first layout as fallback.")
        if len(prs.slide_layouts) > 0:
            title_and_content_layout = prs.slide_layouts[0]
        else:
            raise ValueError("Presentation has no slide layouts to use.")

    # --- Create slides ---
    for i, item in enumerate(slide_data):
        print(f"\n--- Processing Slide {i+1} ---")
        slide = prs.slides.add_slide(title_and_content_layout)
        
        title_shape = slide.shapes.title
        body_shape = None

        if title_shape:
            title_text = item.get("title", "No Title Provided")
            print(f"Setting title: '{title_text}'")
            title_shape.text = title_text
        else:
            print("Warning: Slide layout has no title placeholder.")

        # Find the body placeholder
        for shape in slide.placeholders:
            # We are looking for the main content area, which can be BODY or OBJECT
            if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                if shape.has_text_frame:
                    body_shape = shape
                    print(f"Found body placeholder with type: {shape.placeholder_format.type}")
                    break
        
        if body_shape:
            tf = body_shape.text_frame
            tf.clear() # Clear any default text
            
            points = item.get("points", [])
            print(f"Adding points: {points}")
            
            if isinstance(points, list) and points:
                # Set the first point
                p = tf.paragraphs[0]
                p.text = str(points[0])
                
                # Add subsequent points
                for point_text in points[1:]:
                    p = tf.add_paragraph()
                    p.text = str(point_text)
                    p.level = 0
            elif isinstance(points, str): # Handle if LLM returns a single string instead of a list
                p = tf.paragraphs[0]
                p.text = points
            else:
                print("No points to add or points format is incorrect.")
        else:
            print("Warning: Could not find a body placeholder on this slide.")

    print("\n--- Saving Presentation ---")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path

